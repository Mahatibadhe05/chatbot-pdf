import streamlit as st
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import google.generativeai as genai
import os
import json
import uuid

# =========================
# GEMINI API KEY
# =========================

from dotenv import load_dotenv

load_dotenv()

genai.configure(
    api_key=os.getenv("GEMINI_API_KEY")
)
model = genai.GenerativeModel("gemini-2.5-flash")

# =========================
# CHAT STORAGE
# =========================

CHAT_DIR = "chats"
os.makedirs(CHAT_DIR, exist_ok=True)


def save_chat(
    chat_id,
    messages,
    title="New Chat"
):
    with open(
        f"{CHAT_DIR}/{chat_id}.json",
        "w",
        encoding="utf-8"
    ) as f:

        json.dump(
            {
                "title": title,
                "messages": messages
            },
            f,
            indent=4,
            ensure_ascii=False
        )


def load_chat(chat_id):

    path = f"{CHAT_DIR}/{chat_id}.json"

    if os.path.exists(path):

        with open(
            path,
            "r",
            encoding="utf-8"
        ) as f:

            return json.load(f)

    return {
        "title": "New Chat",
        "messages": []
    }


def get_chats():

    chats = []

    for file in os.listdir(CHAT_DIR):

        if file.endswith(".json"):

            chat_id = file[:-5]

            try:

                data = load_chat(chat_id)

                title = data.get(
                    "title",
                    "New Chat"
                )

            except:

                title = "New Chat"

            chats.append(
                {
                    "id": chat_id,
                    "title": title
                }
            )

    chats.reverse()

    return chats

    
# =========================
# PAGE SETTINGS
# =========================

st.set_page_config(
    page_title="AI PDF Chatbot",
    page_icon="🤖",
    layout="wide"
)

st.markdown("""
<style>
[data-testid="stSidebar"] {
    background-color: #f5f5f5;
}

.stButton button {
    border-radius: 12px;
}

section[data-testid="stFileUploader"] {
    padding: 0px;
    border: none;
}
</style>
""", unsafe_allow_html=True)

# =========================
# SESSION STATE
# =========================

# =========================
# SESSION STATE
# =========================

if "chat_id" not in st.session_state:
    st.session_state.chat_id = str(uuid.uuid4())

if "messages" not in st.session_state:
    st.session_state.messages = []

if "pdf_text" not in st.session_state:
    st.session_state.pdf_text = ""

if "chat_title" not in st.session_state:
    st.session_state.chat_title = "New Chat"
# =========================
# SIDEBAR
# =========================

# =========================
# SIDEBAR
# =========================

with st.sidebar:

    st.title("📚 Chats")

    st.markdown("### 📎 Files")

    uploaded_file = st.file_uploader(
        "",
        type=["pdf", "docx", "pptx"],
        label_visibility="collapsed"
    )

    st.divider()

    if st.button(
        "➕ New Chat",
        use_container_width=True
):

        st.session_state.chat_id = str(uuid.uuid4())
        st.session_state.messages = []
        st.session_state.pdf_text = ""
        st.session_state.chat_title = "New Chat"

        st.rerun()

    st.divider()

    chats = get_chats()

    for chat in chats:

        col1, col2 = st.columns([5, 1])

        with col1:

            if st.button(
                f"📄 {chat['title']}",
                key=f"chat_{chat['id']}",
                use_container_width=True
            ):

                st.session_state.chat_id = chat["id"]

                data = load_chat(chat["id"])

                st.session_state.messages = data["messages"]

                st.session_state.chat_title = data["title"]

                st.rerun()

        with col2:

            if st.button(
                "🗑️",
                key=f"delete_{chat['id']}"
            ):

                os.remove(
                    f"{CHAT_DIR}/{chat['id']}.json"
                )

                st.rerun()
# =========================
# MAIN PAGE
# =========================


# =========================
# PDF UPLOAD
# =========================

# =========================
# FILE UPLOAD
# =========================
if uploaded_file:

    filename = os.path.splitext(
        uploaded_file.name
    )[0]

    st.session_state.chat_title = filename

    text = ""

    # PDF
    if uploaded_file.name.endswith(".pdf"):

        reader = PdfReader(uploaded_file)

        for page in reader.pages:

            extracted = page.extract_text()

            if extracted:
                text += extracted

    # DOCX
    elif uploaded_file.name.endswith(".docx"):

        doc = Document(uploaded_file)

        for para in doc.paragraphs:

            text += para.text + "\n"

    # PPTX
    elif uploaded_file.name.endswith(".pptx"):

        prs = Presentation(uploaded_file)

        for slide in prs.slides:

            for shape in slide.shapes:

                if hasattr(shape, "text"):

                    text += shape.text + "\n"

    st.session_state.pdf_text = text

    st.success("✅ File uploaded successfully!")
# =========================
# DISPLAY CHAT HISTORY
# =========================
if len(st.session_state.messages) == 0:

    st.markdown(
        """
        <div style="text-align:center; padding:80px;">
            <h1>✨ Mahati AI</h1>
            <p>Upload a PDF, DOCX, or PPTX and start chatting.</p>
        </div>
        """,
        unsafe_allow_html=True
    )

for message in st.session_state.messages:

    with st.chat_message(message["role"]):

        st.markdown(message["content"])

# =========================
# CHAT INPUT
# =========================

question = st.chat_input("Ask anything...")

if question:

    st.session_state.messages.append(
        {
            "role": "user",
            "content": question
        }
    )

    save_chat(
        st.session_state.chat_id,
        st.session_state.messages,
        st.session_state.chat_title
    )

    with st.chat_message("user"):
        st.markdown(question)

    try:

        history = ""

        for msg in st.session_state.messages:

            history += (
                f"{msg['role']}: "
                f"{msg['content']}\n"
            )
        prompt = f"""
You are Mahati AI.

When a document is uploaded:

- Answer questions about the uploaded document.
- Respond naturally to conversational messages, greetings,
  compliments, thanks, and casual chat.
- If a question is clearly seeking factual information that is
  unrelated to the uploaded document, politely say:

  "That question is not related to the uploaded document.
  Please ask about the document or start a new chat."

Examples:

User: Hello
Assistant: Hello! How can I help?

User: Good work
Assistant: Thank you!

User: Thanks
Assistant: You're welcome!

User: Who organized this event?
Assistant: (Answer from document)

User: Who is Elon Musk's child?
Assistant: That question is not related to the uploaded document.

DOCUMENT:
{st.session_state.pdf_text[:30000]}

CHAT HISTORY:
{history}

USER:
{question}
"""


        with st.chat_message("assistant"):

            placeholder = st.empty()

            placeholder.markdown(
                "🤔 **Mahati AI is thinking...**"
            )

            full_response = ""

            response = model.generate_content(
                prompt,
                stream=True
            )

            for chunk in response:

                if hasattr(chunk, "text"):

                    full_response += chunk.text

                    placeholder.markdown(
                        full_response + "▌"
                    )

            placeholder.markdown(full_response)

            answer = full_response

        st.session_state.messages.append(
            {
                "role": "assistant",
                "content": answer
            }
        )

        save_chat(
            st.session_state.chat_id,
            st.session_state.messages,
            st.session_state.chat_title
)

    except Exception as e:

        with st.chat_message("assistant"):

            st.error(f"Error: {e}")